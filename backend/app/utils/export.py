"""Document Export Service"""
import os
from typing import Optional
from uuid import UUID
import uuid as uuid_module
from datetime import datetime

from docx import Document as DocxDocument
from docx.shared import Pt, RGBColor, Inches
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
from pptx import Presentation
from pptx.util import Inches, Pt as PptPt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor as PptRGBColor


class WordExporter:
    """Export to Word (.docx) format"""
    
    @staticmethod
    def create_document(
        project_title: str,
        sections: list,
        style_config: dict = None
    ) -> bytes:
        """Create Word document from sections"""
        doc = DocxDocument()
        
        # Add title
        title = doc.add_heading(project_title, level=0)
        title.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER
        
        # Add metadata
        style_config = style_config or {}
        metadata_para = doc.add_paragraph()
        metadata_para.add_run(f"Generated: {datetime.utcnow().strftime('%Y-%m-%d %H:%M:%S')}").italic = True
        
        # Add table of contents
        doc.add_paragraph("Table of Contents", style='Heading 1')
        toc_para = doc.add_paragraph()
        for i, section in enumerate(sections, 1):
            toc_para.add_run(f"{i}. {section.get('title', 'Untitled')}\n")
        
        # Add page break
        doc.add_page_break()
        
        # Add sections
        for i, section in enumerate(sections, 1):
            # Section heading
            heading = doc.add_heading(f"{i}. {section.get('title', 'Untitled')}", level=1)
            
            # Section content
            content = section.get('content', '')
            if content:
                # Parse content for formatting
                paragraphs = content.split('\n')
                for para_text in paragraphs:
                    if para_text.strip():
                        if para_text.strip().startswith('- ') or para_text.strip().startswith('• '):
                            # Bullet point
                            para = doc.add_paragraph(para_text.strip()[2:], style='List Bullet')
                        else:
                            # Regular paragraph
                            para = doc.add_paragraph(para_text.strip())
                            para.alignment = WD_PARAGRAPH_ALIGNMENT.LEFT
            
            # Add spacing
            if i < len(sections):
                doc.add_paragraph()
                doc.add_page_break()
        
        # Save to bytes
        from io import BytesIO
        byte_stream = BytesIO()
        doc.save(byte_stream)
        byte_stream.seek(0)
        return byte_stream.getvalue()
    
    @staticmethod
    def add_formatting(paragraph, style_config: dict):
        """Apply formatting to paragraph"""
        font_name = style_config.get("font", "Calibri")
        font_size = style_config.get("font_size", 12)
        
        for run in paragraph.runs:
            run.font.name = font_name
            run.font.size = Pt(font_size)


class PowerPointExporter:
    """Export to PowerPoint (.pptx) format"""
    
    @staticmethod
    def create_presentation(
        project_title: str,
        sections: list,
        slide_titles: list = None,
        style_config: dict = None
    ) -> bytes:
        """Create PowerPoint presentation from sections"""
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        style_config = style_config or {}
        
        # Title slide
        title_slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        title_shape = title_slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(2))
        title_frame = title_shape.text_frame
        title_frame.word_wrap = True
        p = title_frame.paragraphs[0]
        p.text = project_title
        p.font.size = PptPt(54)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        
        # Subtitle
        subtitle_shape = title_slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(1))
        subtitle_frame = subtitle_shape.text_frame
        p = subtitle_frame.paragraphs[0]
        p.text = f"Generated: {datetime.utcnow().strftime('%B %d, %Y')}"
        p.font.size = PptPt(24)
        p.alignment = PP_ALIGN.CENTER
        
        # Content slides
        for idx, section in enumerate(sections):
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
            
            # Slide title
            slide_title = slide_titles[idx] if slide_titles and idx < len(slide_titles) else section.get('title', f'Slide {idx + 1}')
            title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
            title_frame = title_shape.text_frame
            p = title_frame.paragraphs[0]
            p.text = slide_title
            p.font.size = PptPt(40)
            p.font.bold = True
            
            # Content
            content_shape = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
            content_frame = content_shape.text_frame
            content_frame.word_wrap = True
            
            content = section.get('content', '')
            if content:
                lines = content.split('\n')
                for line in lines:
                    line = line.strip()
                    if line:
                        if line.startswith('- ') or line.startswith('• '):
                            p = content_frame.add_paragraph()
                            p.text = line[2:]
                            p.level = 0
                            p.font.size = PptPt(20)
                        elif line.startswith('  - ') or line.startswith('  • '):
                            p = content_frame.add_paragraph()
                            p.text = line[4:]
                            p.level = 1
                            p.font.size = PptPt(18)
                        else:
                            p = content_frame.add_paragraph()
                            p.text = line
                            p.level = 0
                            p.font.size = PptPt(20)
        
        # Save to bytes
        from io import BytesIO
        byte_stream = BytesIO()
        prs.save(byte_stream)
        byte_stream.seek(0)
        return byte_stream.getvalue()


class ExportService:
    """Main export service coordinating different formats"""
    
    @staticmethod
    def export_document(
        document_id: UUID,
        sections: list,
        export_format: str,
        export_options: dict = None,
        slide_titles: list = None,
        project_title: str = "Document"
    ) -> bytes:
        """Export document in specified format"""
        export_options = export_options or {}
        
        if export_format.lower() == "docx":
            return WordExporter.create_document(project_title, sections, export_options)
        elif export_format.lower() == "pptx":
            return PowerPointExporter.create_presentation(
                project_title, sections, slide_titles, export_options
            )
        else:
            raise ValueError(f"Unsupported export format: {export_format}")
    
    @staticmethod
    def save_export(
        file_bytes: bytes,
        document_id: UUID,
        export_format: str,
        temp_dir: str = "./exports"
    ) -> str:
        """Save exported file to disk and return path"""
        os.makedirs(temp_dir, exist_ok=True)
        
        filename = f"{document_id}_{datetime.utcnow().timestamp()}.{export_format}"
        filepath = os.path.join(temp_dir, filename)
        
        with open(filepath, 'wb') as f:
            f.write(file_bytes)
        
        return filepath


class TemplateService:
    """Service for AI-generated templates (bonus feature)"""
    
    @staticmethod
    def generate_outline_template(
        topic: str,
        document_type: str,
        num_sections: int,
        llm_client,
        style: str = "professional"
    ):
        """Generate AI-suggested outline"""
        from app.integrations import PromptManager
        import json
        import asyncio
        
        prompt = PromptManager.build_outline_prompt(
            topic, document_type, num_sections, style
        )
        
        # Get response from LLM
        loop = asyncio.get_event_loop()
        response = loop.run_until_complete(llm_client.generate_content(prompt))
        
        # Parse JSON response
        try:
            import re
            json_match = re.search(r'\[.*\]', response, re.DOTALL)
            if json_match:
                outline = json.loads(json_match.group())
                return outline
        except json.JSONDecodeError:
            pass
        
        # Fallback: return structured outline
        return [
            {"title": f"Section {i+1}", "description": f"Content for section {i+1}"}
            for i in range(num_sections)
        ]
    
    @staticmethod
    def generate_slide_titles_template(
        topic: str,
        num_slides: int,
        llm_client,
        audience: str = "general"
    ):
        """Generate AI-suggested slide titles"""
        from app.integrations import PromptManager
        import json
        import asyncio
        
        prompt = PromptManager.build_slide_title_prompt(topic, num_slides, audience)
        
        # Get response from LLM
        loop = asyncio.get_event_loop()
        response = loop.run_until_complete(llm_client.generate_content(prompt))
        
        # Parse JSON response
        try:
            import re
            json_match = re.search(r'\[.*\]', response, re.DOTALL)
            if json_match:
                titles = json.loads(json_match.group())
                return titles[:num_slides]
        except json.JSONDecodeError:
            pass
        
        # Fallback: return basic slide titles
        return [f"Slide {i+1}: {topic}" for i in range(num_slides)]
